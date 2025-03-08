import PyPDF2
import docx
import pptx

def extract_text(uploaded_file):
    if uploaded_file.type == "application/pdf":
        return extract_pdf_text(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_docx_text(uploaded_file)
    elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        return extract_ppt_text(uploaded_file)
    return ""

def extract_pdf_text(pdf_file):
    text = ""
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    for page in pdf_reader.pages:
        text += page.extract_text() + " "
    return text.strip()

def extract_docx_text(docx_file):
    doc = docx.Document(docx_file)
    return " ".join([para.text for para in doc.paragraphs])

def extract_ppt_text(ppt_file):
    text = ""
    presentation = pptx.Presentation(ppt_file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text.strip()

